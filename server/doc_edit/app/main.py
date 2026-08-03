"""Document-edit service: xlsx / docx / pptx / txt / md / csv / tsv.

Also supports cross-format conversion via POST /v1/documents/convert.

Contract: docs/document-edit-contract.md
Env: DOC_API_TOKEN, MAX_UPLOAD_MB, TEMP_DIR
"""

from __future__ import annotations

import csv
import io
import json
import os
import re
import shutil
import tempfile
import uuid
from pathlib import Path
from typing import Any

from docx import Document
from fastapi import Depends, FastAPI, File, Form, Header, HTTPException, UploadFile
from fastapi.responses import FileResponse, JSONResponse
from openpyxl import Workbook, load_workbook
from openpyxl.utils import coordinate_to_tuple
from pptx import Presentation
from starlette.background import BackgroundTask

APP_VERSION = "0.4.0"
MAX_OPS = 200
MAX_CELLS_PER_OP = 5000
MAX_SET_TEXT_CHARS = 2 * 1024 * 1024
A1_RE = re.compile(r"^[A-Z]{1,3}[1-9][0-9]{0,6}$")
SUPPORTED = {".xlsx", ".docx", ".pptx", ".txt", ".md", ".csv", ".tsv"}
FORMATS = ["xlsx", "docx", "pptx", "txt", "md", "csv", "tsv"]
# Source format → allowed target formats (same-format always allowed).
CONVERSIONS: dict[str, set[str]] = {
    "txt": {"txt", "md", "docx", "csv", "tsv"},
    "md": {"md", "txt", "docx"},
    "csv": {"csv", "tsv", "txt", "md", "xlsx"},
    "tsv": {"tsv", "csv", "txt", "md", "xlsx"},
    "xlsx": {"xlsx", "csv", "tsv", "txt", "md"},
    "docx": {"docx", "txt", "md"},
    "pptx": {"pptx", "txt", "md"},
}
MEDIA_TYPES = {
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".txt": "text/plain; charset=utf-8",
    ".md": "text/markdown; charset=utf-8",
    ".csv": "text/csv; charset=utf-8",
    ".tsv": "text/tab-separated-values; charset=utf-8",
}

app = FastAPI(title="Expert Chat Document Edit", version=APP_VERSION)


def _token() -> str:
    return os.environ.get("DOC_API_TOKEN", "").strip()


def _max_upload_bytes() -> int:
    mb = float(os.environ.get("MAX_UPLOAD_MB", "20"))
    return int(mb * 1024 * 1024)


def _temp_root() -> Path:
    raw = os.environ.get("TEMP_DIR", "").strip()
    p = Path(raw) if raw else Path(tempfile.gettempdir()) / "expert_chat_doc_edit"
    p.mkdir(parents=True, exist_ok=True)
    return p


def require_auth(authorization: str | None = Header(default=None)) -> None:
    expected = _token()
    if not expected:
        raise HTTPException(
            status_code=500,
            detail=_err("internal", "服务器未配置 DOC_API_TOKEN"),
        )
    if not authorization or not authorization.lower().startswith("bearer "):
        raise HTTPException(
            status_code=401,
            detail=_err("unauthorized", "缺少 Authorization: Bearer"),
        )
    if authorization[7:].strip() != expected:
        raise HTTPException(
            status_code=401,
            detail=_err("unauthorized", "Token 无效"),
        )


def _err(code: str, message: str, details: dict[str, Any] | None = None) -> dict:
    return {"error": {"code": code, "message": message, "details": details or {}}}


@app.get("/v1/health")
def health() -> dict[str, Any]:
    return {
        "ok": True,
        "version": APP_VERSION,
        "formats": FORMATS,
        "conversions": {k: sorted(v) for k, v in CONVERSIONS.items()},
    }


@app.post("/v1/documents/edit")
async def edit_document(
    file: UploadFile = File(...),
    patch: str = Form(...),
    filename: str | None = Form(default=None),
    _: None = Depends(require_auth),
) -> FileResponse:
    job_id = uuid.uuid4().hex
    job_dir = _temp_root() / job_id
    job_dir.mkdir(parents=True, exist_ok=True)

    def cleanup() -> None:
        shutil.rmtree(job_dir, ignore_errors=True)

    try:
        raw_name = filename or file.filename or "input.bin"
        lower = raw_name.lower()
        ext = Path(lower).suffix
        if ext not in SUPPORTED:
            raise HTTPException(
                status_code=400,
                detail=_err("unsupported_format", f"仅支持 {', '.join(sorted(SUPPORTED))}"),
            )

        data = await file.read()
        if len(data) > _max_upload_bytes():
            raise HTTPException(
                status_code=400,
                detail=_err(
                    "file_too_large",
                    f"文件超过 {_max_upload_bytes() // (1024 * 1024)} MB 限制",
                ),
            )
        if not data:
            raise HTTPException(
                status_code=400,
                detail=_err("patch_invalid", "上传文件为空"),
            )

        try:
            patch_obj = json.loads(patch)
        except json.JSONDecodeError as e:
            raise HTTPException(
                status_code=400,
                detail=_err("patch_invalid", f"patch 不是合法 JSON：{e}"),
            ) from e

        try:
            validated = _validate_patch(patch_obj, ext)
        except ValueError as e:
            raise HTTPException(
                status_code=400,
                detail=_err("patch_invalid", str(e)),
            ) from e

        in_path = job_dir / f"input{ext}"
        out_path = job_dir / f"output{ext}"
        in_path.write_bytes(data)

        try:
            if ext == ".xlsx":
                _apply_xlsx(in_path, out_path, validated)
            elif ext == ".docx":
                _apply_docx(in_path, out_path, validated)
            elif ext == ".pptx":
                _apply_pptx(in_path, out_path, validated)
            else:
                _apply_text(in_path, out_path, validated)
        except ValueError as e:
            raise HTTPException(
                status_code=400,
                detail=_err("patch_invalid", str(e)),
            ) from e
        except Exception as e:  # noqa: BLE001
            raise HTTPException(
                status_code=500,
                detail=_err("internal", f"写入失败：{e}"),
            ) from e

        out_name = validated.get("output_filename") or _default_out_name(raw_name, ext)
        return FileResponse(
            path=out_path,
            filename=out_name,
            media_type=MEDIA_TYPES[ext],
            background=BackgroundTask(cleanup),
        )
    except HTTPException:
        cleanup()
        raise
    except Exception as e:  # noqa: BLE001
        cleanup()
        raise HTTPException(
            status_code=500,
            detail=_err("internal", f"内部错误：{e}"),
        ) from e


@app.post("/v1/documents/convert")
async def convert_document(
    file: UploadFile = File(...),
    target_format: str = Form(...),
    filename: str | None = Form(default=None),
    output_filename: str | None = Form(default=None),
    _: None = Depends(require_auth),
) -> FileResponse:
    """Convert an uploaded file to another supported format."""
    job_id = uuid.uuid4().hex
    job_dir = _temp_root() / job_id
    job_dir.mkdir(parents=True, exist_ok=True)

    def cleanup() -> None:
        shutil.rmtree(job_dir, ignore_errors=True)

    try:
        raw_name = filename or file.filename or "input.bin"
        src_ext = Path(raw_name.lower()).suffix
        if src_ext not in SUPPORTED:
            raise HTTPException(
                status_code=400,
                detail=_err("unsupported_format", f"仅支持 {', '.join(sorted(SUPPORTED))}"),
            )
        src_fmt = src_ext.lstrip(".")
        tgt_fmt = str(target_format or "").strip().lower().lstrip(".")
        if tgt_fmt not in FORMATS:
            raise HTTPException(
                status_code=400,
                detail=_err("unsupported_format", f'不支持的 target_format="{target_format}"'),
            )
        allowed = CONVERSIONS.get(src_fmt, set())
        if tgt_fmt not in allowed:
            raise HTTPException(
                status_code=400,
                detail=_err(
                    "unsupported_format",
                    f"不支持 {src_fmt} → {tgt_fmt}。"
                    f"可选：{', '.join(sorted(allowed)) or '无'}",
                ),
            )

        data = await file.read()
        if len(data) > _max_upload_bytes():
            raise HTTPException(
                status_code=400,
                detail=_err(
                    "file_too_large",
                    f"文件超过 {_max_upload_bytes() // (1024 * 1024)} MB 限制",
                ),
            )
        if not data:
            raise HTTPException(
                status_code=400,
                detail=_err("patch_invalid", "上传文件为空"),
            )

        out_name = _sanitize_output_filename(output_filename, raw_name, f".{tgt_fmt}")
        in_path = job_dir / f"input{src_ext}"
        out_path = job_dir / f"output.{tgt_fmt}"
        in_path.write_bytes(data)

        try:
            if src_fmt == tgt_fmt:
                if src_fmt in {"txt", "md", "csv", "tsv"}:
                    # Normalize text encodings to UTF-8 on same-format pass-through.
                    _write_text_file(out_path, _read_text_file(in_path))
                else:
                    out_path.write_bytes(data)
            else:
                _convert_file(in_path, src_fmt, out_path, tgt_fmt)
        except ValueError as e:
            raise HTTPException(
                status_code=400,
                detail=_err("unsupported_format", str(e)),
            ) from e
        except Exception as e:  # noqa: BLE001
            raise HTTPException(
                status_code=500,
                detail=_err("internal", f"转换失败：{e}"),
            ) from e

        return FileResponse(
            path=out_path,
            filename=out_name,
            media_type=MEDIA_TYPES[f".{tgt_fmt}"],
            background=BackgroundTask(cleanup),
        )
    except HTTPException:
        cleanup()
        raise
    except Exception as e:  # noqa: BLE001
        cleanup()
        raise HTTPException(
            status_code=500,
            detail=_err("internal", f"内部错误：{e}"),
        ) from e


def _default_out_name(raw: str, ext: str) -> str:
    base = Path(raw).stem or "edited"
    return f"{base}_edited{ext}"


def _sanitize_output_filename(raw: Any, source_name: str, out_ext: str) -> str:
    if raw is not None and str(raw).strip():
        out_s = str(raw).strip()
        if "/" in out_s or "\\" in out_s or ".." in out_s:
            raise HTTPException(
                status_code=400,
                detail=_err("patch_invalid", "output_filename 非法"),
            )
        if len(out_s) > 180:
            raise HTTPException(
                status_code=400,
                detail=_err("patch_invalid", "output_filename 过长"),
            )
        if not out_s.lower().endswith(out_ext.lower()):
            out_s = f"{Path(out_s).stem}{out_ext}"
        return out_s
    base = Path(source_name).stem or "converted"
    return f"{base}{out_ext}"


def _validate_patch(obj: Any, file_ext: str) -> dict[str, Any]:
    if not isinstance(obj, dict):
        raise ValueError("patch 必须是对象")
    if obj.get("schema_version") != 1:
        raise ValueError(f"不支持的 schema_version={obj.get('schema_version')}（需要 1）")
    fmt = str(obj.get("format", "")).lower()
    expected = file_ext.lstrip(".")
    if fmt not in set(FORMATS):
        raise ValueError(f'不支持的 format="{fmt}"')
    if fmt != expected:
        raise ValueError(f'format="{fmt}" 与上传文件扩展名 {file_ext} 不一致')
    ops = obj.get("ops")
    if not isinstance(ops, list) or not ops:
        raise ValueError("ops 必须是非空数组")
    if len(ops) > MAX_OPS:
        raise ValueError(f"ops 超过上限 {MAX_OPS}")

    out = obj.get("output_filename")
    if out is not None:
        out_s = str(out).strip()
        if not out_s or "/" in out_s or "\\" in out_s or ".." in out_s:
            raise ValueError("output_filename 非法")
        if len(out_s) > 180:
            raise ValueError("output_filename 过长")
        obj = {**obj, "output_filename": out_s}

    text_ops = {"replace_text", "set_text"}
    allowed = {
        "xlsx": {"set_cells", "set_range", "add_sheet", "ensure_sheet"},
        "docx": {"replace_text"},
        "pptx": {"set_shape_text"},
        "txt": text_ops,
        "md": text_ops,
        "csv": text_ops,
        "tsv": text_ops,
    }[fmt]

    for i, op in enumerate(ops):
        if not isinstance(op, dict):
            raise ValueError(f"ops[{i}] 必须是对象")
        kind = str(op.get("op", "")).lower()
        if kind not in allowed:
            raise ValueError(f'format={fmt} 不支持 op="{kind}"（ops[{i}]）')
        if kind == "set_cells":
            cells = op.get("cells")
            if not isinstance(cells, dict) or not cells:
                raise ValueError(f"ops[{i}].cells 必须是非空对象")
            if len(cells) > MAX_CELLS_PER_OP:
                raise ValueError(f"ops[{i}].cells 过多")
            for addr in cells:
                if not A1_RE.match(str(addr).strip().upper()):
                    raise ValueError(f"ops[{i}] 非法地址 {addr}")
        elif kind == "set_range":
            start = str(op.get("start", "")).strip().upper()
            if not A1_RE.match(start):
                raise ValueError(f"ops[{i}].start 非法")
            values = op.get("values")
            if not isinstance(values, list) or not values:
                raise ValueError(f"ops[{i}].values 必须是非空二维数组")
            n = sum(len(r) if isinstance(r, list) else 0 for r in values)
            if n > MAX_CELLS_PER_OP:
                raise ValueError(f"ops[{i}].values 单元格过多")
        elif kind in {"add_sheet", "ensure_sheet"}:
            name = str(op.get("name", "")).strip()
            if not name or len(name) > 31:
                raise ValueError(f"ops[{i}].name 非法")
            if re.search(r"[\\/*?:\[\]]", name):
                raise ValueError(f"ops[{i}].name 含非法字符")
        elif kind == "replace_text":
            find = str(op.get("find", ""))
            if not find:
                raise ValueError(f"ops[{i}].find 不能为空")
        elif kind == "set_text":
            if "text" not in op or op.get("text") is None:
                raise ValueError(f"ops[{i}].text 不能为 null")
            text = op.get("text")
            text_s = text if isinstance(text, str) else str(text)
            if len(text_s) > MAX_SET_TEXT_CHARS:
                raise ValueError(f"ops[{i}].text 过长（>{MAX_SET_TEXT_CHARS}）")
        elif kind == "set_shape_text":
            slide = int(op.get("slide", 1))
            if slide < 1:
                raise ValueError(f"ops[{i}].slide 必须 ≥ 1")
            shape = int(op.get("shape", 0))
            if shape < 0:
                raise ValueError(f"ops[{i}].shape 必须 ≥ 0")
    return obj


def _apply_xlsx(in_path: Path, out_path: Path, patch: dict[str, Any]) -> None:
    try:
        wb = load_workbook(in_path)
    except Exception:
        wb = Workbook()

    for op in patch["ops"]:
        kind = str(op["op"]).lower()
        if kind == "add_sheet":
            name = str(op["name"]).strip()
            if name in wb.sheetnames:
                raise ValueError(f'工作表已存在："{name}"')
            wb.create_sheet(name)
        elif kind == "ensure_sheet":
            name = str(op["name"]).strip()
            if name not in wb.sheetnames:
                wb.create_sheet(name)
        elif kind == "set_cells":
            ws = _sheet(wb, op.get("sheet"))
            for addr, value in op["cells"].items():
                ws[str(addr).strip().upper()] = _cell_py(value)
        elif kind == "set_range":
            ws = _sheet(wb, op.get("sheet"))
            start = str(op["start"]).strip().upper()
            row0, col0 = coordinate_to_tuple(start)
            for r, row in enumerate(op["values"]):
                if not isinstance(row, list):
                    raise ValueError("set_range 行必须是数组")
                for c, value in enumerate(row):
                    ws.cell(row=row0 + r, column=col0 + c, value=_cell_py(value))
    wb.save(out_path)


def _apply_docx(in_path: Path, out_path: Path, patch: dict[str, Any]) -> None:
    doc = Document(str(in_path))
    for op in patch["ops"]:
        if str(op["op"]).lower() != "replace_text":
            continue
        find = str(op["find"])
        replace = str(op.get("replace", ""))
        all_occ = bool(op.get("all", True))
        _docx_replace(doc, find, replace, all_occ)
    doc.save(str(out_path))


def _docx_replace(doc: Document, find: str, replace: str, all_occ: bool) -> None:
    replaced = 0

    def do_para(p: Any) -> None:
        nonlocal replaced
        if find not in p.text:
            return
        if all_occ or replaced == 0:
            # Simple full-paragraph text rewrite keeps runs minimal but robust.
            new_text = p.text.replace(find, replace) if all_occ else p.text.replace(find, replace, 1)
            if new_text != p.text:
                for run in p.runs:
                    run.text = ""
                if p.runs:
                    p.runs[0].text = new_text
                else:
                    p.add_run(new_text)
                replaced += 1 if not all_occ else p.text.count(replace)  # best-effort

    for p in doc.paragraphs:
        do_para(p)
        if not all_occ and replaced:
            break
    if all_occ or not replaced:
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        do_para(p)
                        if not all_occ and replaced:
                            return


def _read_text_file(path: Path) -> str:
    raw = path.read_bytes()
    if not raw:
        return ""
    # Strip UTF-8 BOM if present.
    if raw.startswith(b"\xef\xbb\xbf"):
        raw = raw[3:]
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        # Best-effort for Chinese Windows exports.
        try:
            return raw.decode("gb18030")
        except UnicodeDecodeError:
            return raw.decode("utf-8", errors="replace")


def _write_text_file(path: Path, text: str) -> None:
    path.write_bytes(text.encode("utf-8"))


def _apply_text(in_path: Path, out_path: Path, patch: dict[str, Any]) -> None:
    content = _read_text_file(in_path)
    for op in patch["ops"]:
        kind = str(op["op"]).lower()
        if kind == "set_text":
            text = op.get("text")
            content = text if isinstance(text, str) else str(text)
        elif kind == "replace_text":
            find = str(op["find"])
            replace = str(op.get("replace", ""))
            all_occ = bool(op.get("all", True))
            if all_occ:
                content = content.replace(find, replace)
            else:
                content = content.replace(find, replace, 1)
    _write_text_file(out_path, content)


def _apply_pptx(in_path: Path, out_path: Path, patch: dict[str, Any]) -> None:
    prs = Presentation(str(in_path))
    for op in patch["ops"]:
        if str(op["op"]).lower() != "set_shape_text":
            continue
        slide_i = int(op.get("slide", 1)) - 1
        shape_i = int(op.get("shape", 0))
        text = str(op.get("text", ""))
        if slide_i < 0 or slide_i >= len(prs.slides):
            raise ValueError(f"幻灯片不存在：slide={slide_i + 1}")
        slide = prs.slides[slide_i]
        text_shapes = [s for s in slide.shapes if getattr(s, "has_text_frame", False)]
        if shape_i < 0 or shape_i >= len(text_shapes):
            raise ValueError(
                f"幻灯片 {slide_i + 1} 没有 shape 索引 {shape_i}（共 {len(text_shapes)} 个文本框）"
            )
        shape = text_shapes[shape_i]
        tf = shape.text_frame
        if tf.paragraphs:
            # Clear then set first paragraph.
            for pi, para in enumerate(tf.paragraphs):
                if pi == 0:
                    if para.runs:
                        para.runs[0].text = text
                        for run in para.runs[1:]:
                            run.text = ""
                    else:
                        para.text = text
                else:
                    para.text = ""
        else:
            tf.text = text
    prs.save(str(out_path))


def _sheet(wb: Any, name: Any) -> Any:
    if name is None or str(name).strip() == "":
        return wb[wb.sheetnames[0]]
    n = str(name).strip()
    if n not in wb.sheetnames:
        raise ValueError(f'找不到工作表："{n}"')
    return wb[n]


def _cell_py(value: Any) -> Any:
    if value is None or isinstance(value, (bool, int, float, str)):
        return value
    raise ValueError("单元格值类型无效")


def _convert_file(in_path: Path, src_fmt: str, out_path: Path, tgt_fmt: str) -> None:
    """Dispatch cross-format conversion."""
    if src_fmt in {"txt", "md"} and tgt_fmt in {"txt", "md"}:
        _write_text_file(out_path, _read_text_file(in_path))
        return
    if src_fmt in {"txt", "md"} and tgt_fmt == "docx":
        _text_to_docx(_read_text_file(in_path), out_path)
        return
    if src_fmt == "txt" and tgt_fmt in {"csv", "tsv"}:
        _lines_to_delimited(_read_text_file(in_path), out_path, tgt_fmt)
        return
    if src_fmt in {"csv", "tsv"} and tgt_fmt in {"csv", "tsv", "txt", "md"}:
        rows = _read_table(in_path, src_fmt)
        if tgt_fmt in {"csv", "tsv"}:
            _write_table(out_path, rows, tgt_fmt)
        else:
            _write_text_file(out_path, _table_to_plain(rows, markdown=tgt_fmt == "md"))
        return
    if src_fmt in {"csv", "tsv"} and tgt_fmt == "xlsx":
        _table_to_xlsx(_read_table(in_path, src_fmt), out_path)
        return
    if src_fmt == "xlsx" and tgt_fmt in {"csv", "tsv", "txt", "md"}:
        rows = _xlsx_to_rows(in_path)
        if tgt_fmt in {"csv", "tsv"}:
            _write_table(out_path, rows, tgt_fmt)
        else:
            _write_text_file(out_path, _table_to_plain(rows, markdown=tgt_fmt == "md"))
        return
    if src_fmt == "docx" and tgt_fmt in {"txt", "md"}:
        _write_text_file(out_path, _docx_to_text(in_path))
        return
    if src_fmt == "pptx" and tgt_fmt in {"txt", "md"}:
        _write_text_file(out_path, _pptx_to_text(in_path))
        return
    raise ValueError(f"未实现的转换：{src_fmt} → {tgt_fmt}")


def _text_to_docx(text: str, out_path: Path) -> None:
    doc = Document()
    # Preserve blank lines as empty paragraphs.
    lines = text.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    if not lines:
        doc.add_paragraph("")
    else:
        for line in lines:
            doc.add_paragraph(line)
    doc.save(str(out_path))


def _lines_to_delimited(text: str, out_path: Path, fmt: str) -> None:
    rows = [[line] for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")]
    _write_table(out_path, rows, fmt)


def _read_table(path: Path, fmt: str) -> list[list[str]]:
    text = _read_text_file(path)
    delim = "\t" if fmt == "tsv" else ","
    reader = csv.reader(io.StringIO(text), delimiter=delim)
    return [list(row) for row in reader]


def _write_table(path: Path, rows: list[list[Any]], fmt: str) -> None:
    delim = "\t" if fmt == "tsv" else ","
    buf = io.StringIO()
    writer = csv.writer(buf, delimiter=delim, lineterminator="\n")
    for row in rows:
        writer.writerow(["" if c is None else c for c in row])
    _write_text_file(path, buf.getvalue())


def _table_to_plain(rows: list[list[Any]], *, markdown: bool) -> str:
    if not rows:
        return ""
    str_rows = [["" if c is None else str(c) for c in row] for row in rows]
    if not markdown:
        return "\n".join("\t".join(r) for r in str_rows)
    # Simple GFM table: first row as header when ≥2 rows.
    if len(str_rows) == 1:
        return "| " + " | ".join(str_rows[0]) + " |"
    header = str_rows[0]
    width = max(len(r) for r in str_rows)
    header = header + [""] * (width - len(header))
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    for row in str_rows[1:]:
        padded = row + [""] * (width - len(row))
        lines.append("| " + " | ".join(padded) + " |")
    return "\n".join(lines)


def _table_to_xlsx(rows: list[list[Any]], out_path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    for r_i, row in enumerate(rows, start=1):
        for c_i, value in enumerate(row, start=1):
            ws.cell(row=r_i, column=c_i, value=value)
    wb.save(out_path)


def _xlsx_to_rows(path: Path) -> list[list[Any]]:
    try:
        wb = load_workbook(path, data_only=True)
    except Exception as e:  # noqa: BLE001
        raise ValueError(f"无法读取 xlsx：{e}") from e
    ws = wb[wb.sheetnames[0]]
    rows: list[list[Any]] = []
    for row in ws.iter_rows(values_only=True):
        # Drop fully empty trailing-style rows later; keep structure.
        cells = list(row)
        if all(c is None or str(c).strip() == "" for c in cells):
            continue
        rows.append(["" if c is None else c for c in cells])
    return rows


def _docx_to_text(path: Path) -> str:
    doc = Document(str(path))
    parts: list[str] = []
    for p in doc.paragraphs:
        parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _pptx_to_text(path: Path) -> str:
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
        parts.append("")
    return "\n".join(parts).strip() + ("\n" if parts else "")


@app.exception_handler(HTTPException)
async def http_exception_handler(_request: Any, exc: HTTPException) -> JSONResponse:
    detail = exc.detail
    if isinstance(detail, dict) and "error" in detail:
        return JSONResponse(status_code=exc.status_code, content=detail)
    return JSONResponse(status_code=exc.status_code, content=_err("internal", str(detail)))
