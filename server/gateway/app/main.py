from __future__ import annotations

import asyncio
import json
import os
import re
import sqlite3
import time
import uuid
from collections import defaultdict, deque
from contextlib import asynccontextmanager, contextmanager
from datetime import datetime, timedelta, timezone
from pathlib import Path
from threading import Lock
from typing import Any

import httpx
from docx import Document
from fastapi import Depends, FastAPI, File, Header, HTTPException, Query, Request, UploadFile
from fastapi.responses import FileResponse, JSONResponse
from openpyxl import load_workbook
from pydantic import BaseModel, Field
from pptx import Presentation
from pypdf import PdfReader

from doc_edit.app.main import (
    CONVERSIONS as DOCUMENT_CONVERSIONS,
    FORMATS as DOCUMENT_FORMATS,
    require_auth as document_require_auth,
    router as document_router,
)

from .auth import GatewayAuthenticator, Principal
from .module_registry import (
    GatewayCapability,
    GatewayModule,
    GatewayModuleRegistry,
)


def _now() -> str:
    return datetime.now(timezone.utc).isoformat()


DATA_DIR = Path(os.getenv("GATEWAY_DATA_DIR", "./data")).resolve()
UPLOAD_DIR = DATA_DIR / "uploads"
DB_PATH = DATA_DIR / "gateway.sqlite3"
LLM_BASE_URL = os.getenv("LLM_BASE_URL", "").strip().rstrip("/")
LLM_API_KEY = os.getenv("LLM_API_KEY", "").strip()
LLM_MODEL = os.getenv("LLM_MODEL", "").strip()
MAX_FILE_BYTES = int(os.getenv("GATEWAY_MAX_FILE_MB", "50")) * 1024 * 1024
CHUNK_CHARS = max(2000, int(os.getenv("GATEWAY_CHUNK_CHARS", "12000")))
CONCURRENCY = max(1, int(os.getenv("GATEWAY_CONCURRENCY", "2")))
TASK_RETENTION_DAYS = max(0, int(os.getenv("GATEWAY_TASK_RETENTION_DAYS", "30")))
CLEANUP_INTERVAL_SECONDS = max(
    60, int(os.getenv("GATEWAY_CLEANUP_INTERVAL_SECONDS", str(6 * 60 * 60)))
)
DEFAULT_CONCURRENT_TASKS = max(
    1, int(os.getenv("GATEWAY_DEFAULT_CONCURRENT_TASKS", "2"))
)
DEFAULT_STORAGE_BYTES = max(
    MAX_FILE_BYTES,
    int(os.getenv("GATEWAY_DEFAULT_STORAGE_MB", "512")) * 1024 * 1024,
)
RATE_LIMIT_REQUESTS = max(10, int(os.getenv("GATEWAY_RATE_LIMIT_REQUESTS", "180")))
RATE_LIMIT_UPLOADS = max(1, int(os.getenv("GATEWAY_RATE_LIMIT_UPLOADS", "20")))

PERMISSION_LABELS = {
    "gateway.use": "使用 Gateway",
    "files.write": "上传和删除文件",
    "tasks.create": "创建长任务",
    "tasks.read": "查看自己的长任务",
    "tasks.cancel": "取消自己的长任务",
    "documents.edit": "编辑文档",
    "documents.convert": "转换文档格式",
}
DEFAULT_PERMISSIONS = frozenset(
    item.strip()
    for item in os.getenv(
        "GATEWAY_DEFAULT_PERMISSIONS", ",".join(PERMISSION_LABELS)
    ).split(",")
    if item.strip() in PERMISSION_LABELS
)
CAPABILITY_PERMISSIONS = {
    "long_tasks": frozenset({"files.write", "tasks.create", "tasks.read"}),
    "document_edit": frozenset({"documents.edit"}),
    "document_convert": frozenset({"documents.convert"}),
}

AUTH = GatewayAuthenticator()

MODULES = GatewayModuleRegistry(
    (
        GatewayModule(
            name="long_tasks",
            capabilities=(
                GatewayCapability(
                    id="long_tasks",
                    metadata={
                        "file_limit_bytes": MAX_FILE_BYTES,
                        "durable_events": True,
                        "idempotent_create": True,
                    },
                ),
            ),
        ),
        GatewayModule(
            name="documents",
            router=document_router,
            capabilities=(
                GatewayCapability(
                    id="document_edit",
                    metadata={
                        "formats": DOCUMENT_FORMATS,
                        "patch_schema_versions": [1],
                    },
                ),
                GatewayCapability(
                    id="document_convert",
                    metadata={
                        "conversions": {
                            key: sorted(value)
                            for key, value in DOCUMENT_CONVERSIONS.items()
                        },
                    },
                ),
            ),
        ),
    )
)

_running: dict[str, asyncio.Task[None]] = {}
_semaphore = asyncio.Semaphore(CONCURRENCY)
_shutting_down = False
_maintenance_task: asyncio.Task[None] | None = None
_rate_windows: dict[tuple[str, str], deque[float]] = defaultdict(deque)
_rate_lock = Lock()


def _connect() -> sqlite3.Connection:
    connection = sqlite3.connect(DB_PATH, timeout=30)
    connection.row_factory = sqlite3.Row
    connection.execute("PRAGMA journal_mode=WAL")
    connection.execute("PRAGMA foreign_keys=ON")
    return connection


@contextmanager
def _db():
    connection = _connect()
    try:
        with connection:
            yield connection
    finally:
        connection.close()


def _init_db() -> None:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    with _db() as db:
        db.executescript(
            """
            CREATE TABLE IF NOT EXISTS files (
                id TEXT PRIMARY KEY,
                owner_sub TEXT NOT NULL DEFAULT 'legacy-owner',
                name TEXT NOT NULL,
                mime_type TEXT NOT NULL,
                size_bytes INTEGER NOT NULL,
                path TEXT NOT NULL,
                created_at TEXT NOT NULL
            );
            CREATE TABLE IF NOT EXISTS tasks (
                id TEXT PRIMARY KEY,
                owner_sub TEXT NOT NULL DEFAULT 'legacy-owner',
                client_request_id TEXT,
                status TEXT NOT NULL,
                prompt TEXT NOT NULL,
                instructions TEXT NOT NULL,
                messages_json TEXT NOT NULL,
                file_ids_json TEXT NOT NULL,
                model TEXT NOT NULL,
                output_text TEXT NOT NULL DEFAULT '',
                progress REAL NOT NULL DEFAULT 0,
                detail TEXT,
                error TEXT,
                cancel_requested INTEGER NOT NULL DEFAULT 0,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL
            );
            CREATE TABLE IF NOT EXISTS events (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                task_id TEXT NOT NULL,
                type TEXT NOT NULL,
                data_json TEXT NOT NULL,
                created_at TEXT NOT NULL,
                FOREIGN KEY(task_id) REFERENCES tasks(id) ON DELETE CASCADE
            );
            CREATE INDEX IF NOT EXISTS idx_events_task_id_id
                ON events(task_id, id);
            CREATE TABLE IF NOT EXISTS user_entitlements (
                owner_sub TEXT PRIMARY KEY,
                display_name TEXT NOT NULL,
                enabled INTEGER NOT NULL DEFAULT 1,
                is_admin INTEGER NOT NULL DEFAULT 0,
                permissions_json TEXT NOT NULL,
                max_concurrent_tasks INTEGER NOT NULL,
                storage_quota_bytes INTEGER NOT NULL,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL
            );
            CREATE TABLE IF NOT EXISTS audit_logs (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                actor_sub TEXT NOT NULL,
                action TEXT NOT NULL,
                target_sub TEXT,
                resource_type TEXT,
                resource_id TEXT,
                detail_json TEXT NOT NULL DEFAULT '{}',
                created_at TEXT NOT NULL
            );
            """
        )
        task_columns = {
            row["name"] for row in db.execute("PRAGMA table_info(tasks)").fetchall()
        }
        if "client_request_id" not in task_columns:
            db.execute("ALTER TABLE tasks ADD COLUMN client_request_id TEXT")
        if "owner_sub" not in task_columns:
            db.execute(
                "ALTER TABLE tasks ADD COLUMN owner_sub TEXT NOT NULL "
                "DEFAULT 'legacy-owner'"
            )
            db.execute(
                "UPDATE tasks SET owner_sub = ? WHERE owner_sub = 'legacy-owner'",
                (AUTH.legacy_owner,),
            )
        file_columns = {
            row["name"] for row in db.execute("PRAGMA table_info(files)").fetchall()
        }
        if "owner_sub" not in file_columns:
            db.execute(
                "ALTER TABLE files ADD COLUMN owner_sub TEXT NOT NULL "
                "DEFAULT 'legacy-owner'"
            )
            db.execute(
                "UPDATE files SET owner_sub = ? WHERE owner_sub = 'legacy-owner'",
                (AUTH.legacy_owner,),
            )
        db.execute("DROP INDEX IF EXISTS idx_tasks_client_request_id")
        db.execute(
            "CREATE UNIQUE INDEX IF NOT EXISTS idx_tasks_owner_request "
            "ON tasks(owner_sub, client_request_id) "
            "WHERE client_request_id IS NOT NULL"
        )
        db.execute(
            "CREATE INDEX IF NOT EXISTS idx_tasks_status_updated_at "
            "ON tasks(status, updated_at)"
        )
        db.execute(
            "CREATE INDEX IF NOT EXISTS idx_tasks_owner_status_updated "
            "ON tasks(owner_sub, status, updated_at)"
        )
        db.execute(
            "CREATE INDEX IF NOT EXISTS idx_files_owner_created "
            "ON files(owner_sub, created_at)"
        )
        db.execute(
            "CREATE INDEX IF NOT EXISTS idx_audit_created ON audit_logs(created_at)"
        )


def _cleanup_expired_tasks(*, now: datetime | None = None) -> int:
    """Delete terminal tasks after their recovery window has elapsed.

    Events are removed by the foreign-key cascade. Queued and running tasks are
    deliberately excluded so maintenance can never interrupt active work.
    A retention value of zero disables automatic deletion.
    """
    if TASK_RETENTION_DAYS <= 0:
        return 0
    reference = now or datetime.now(timezone.utc)
    if reference.tzinfo is None:
        reference = reference.replace(tzinfo=timezone.utc)
    cutoff = (reference - timedelta(days=TASK_RETENTION_DAYS)).isoformat()
    with _db() as db:
        cursor = db.execute(
            "DELETE FROM tasks "
            "WHERE status IN ('completed', 'failed', 'cancelled') "
            "AND updated_at < ?",
            (cutoff,),
        )
        return max(0, cursor.rowcount)


def _audit(
    actor_sub: str,
    action: str,
    *,
    target_sub: str | None = None,
    resource_type: str | None = None,
    resource_id: str | None = None,
    detail: dict[str, Any] | None = None,
) -> None:
    with _db() as db:
        db.execute(
            """
            INSERT INTO audit_logs(
                actor_sub, action, target_sub, resource_type, resource_id,
                detail_json, created_at
            ) VALUES (?, ?, ?, ?, ?, ?, ?)
            """,
            (
                actor_sub,
                action,
                target_sub,
                resource_type,
                resource_id,
                json.dumps(detail or {}, ensure_ascii=False),
                _now(),
            ),
        )


def _provision_principal(principal: Principal) -> sqlite3.Row:
    now = _now()
    with _db() as db:
        existing = db.execute(
            "SELECT * FROM user_entitlements WHERE owner_sub = ?",
            (principal.subject,),
        ).fetchone()
        if existing is None:
            is_admin = int(principal.bootstrap_admin)
            db.execute(
                """
                INSERT INTO user_entitlements(
                    owner_sub, display_name, enabled, is_admin,
                    permissions_json, max_concurrent_tasks,
                    storage_quota_bytes, created_at, updated_at
                ) VALUES (?, ?, 1, ?, ?, ?, ?, ?, ?)
                """,
                (
                    principal.subject,
                    principal.display_name,
                    is_admin,
                    json.dumps(sorted(DEFAULT_PERMISSIONS)),
                    DEFAULT_CONCURRENT_TASKS,
                    DEFAULT_STORAGE_BYTES,
                    now,
                    now,
                ),
            )
            existing = db.execute(
                "SELECT * FROM user_entitlements WHERE owner_sub = ?",
                (principal.subject,),
            ).fetchone()
        else:
            is_admin = int(bool(existing["is_admin"]) or principal.bootstrap_admin)
            db.execute(
                "UPDATE user_entitlements SET display_name = ?, is_admin = ?, "
                "updated_at = ? WHERE owner_sub = ?",
                (principal.display_name, is_admin, now, principal.subject),
            )
            existing = db.execute(
                "SELECT * FROM user_entitlements WHERE owner_sub = ?",
                (principal.subject,),
            ).fetchone()
    assert existing is not None
    return existing


def _principal(
    authorization: str | None = Header(default=None),
) -> Principal:
    principal = AUTH.authenticate(authorization)
    entitlement = _provision_principal(principal)
    if not bool(entitlement["enabled"]):
        raise HTTPException(status_code=403, detail="该 Gateway 账户已停用。")
    return principal


def _permissions(entitlement: sqlite3.Row) -> frozenset[str]:
    try:
        raw = json.loads(entitlement["permissions_json"])
    except (TypeError, json.JSONDecodeError):
        raw = []
    return frozenset(str(item) for item in raw if str(item) in PERMISSION_LABELS)


def _consume_rate_limit(subject: str, group: str) -> None:
    limit = RATE_LIMIT_UPLOADS if group == "upload" else RATE_LIMIT_REQUESTS
    now = time.monotonic()
    cutoff = now - 60.0
    key = (subject, group)
    with _rate_lock:
        window = _rate_windows[key]
        while window and window[0] < cutoff:
            window.popleft()
        if len(window) >= limit:
            retry_after = max(1, int(60 - (now - window[0])))
            raise HTTPException(
                status_code=429,
                detail="请求过于频繁，请稍后重试。",
                headers={"Retry-After": str(retry_after)},
            )
        window.append(now)


def _authorize(principal: Principal, permission: str, *, group: str = "api") -> sqlite3.Row:
    entitlement = _provision_principal(principal)
    if not bool(entitlement["enabled"]):
        raise HTTPException(status_code=403, detail="该 Gateway 账户已停用。")
    if not bool(entitlement["is_admin"]) and permission not in _permissions(entitlement):
        raise HTTPException(status_code=403, detail=f"账户缺少权限：{permission}")
    _consume_rate_limit(principal.subject, group)
    return entitlement


def _require(permission: str, *, group: str = "api"):
    def dependency(principal: Principal = Depends(_principal)) -> Principal:
        _authorize(principal, permission, group=group)
        return principal

    return dependency


def _require_admin(principal: Principal = Depends(_principal)) -> Principal:
    entitlement = _provision_principal(principal)
    if not bool(entitlement["is_admin"]):
        raise HTTPException(status_code=403, detail="需要 Gateway 管理员权限。")
    _consume_rate_limit(principal.subject, "api")
    return principal


def _document_auth(
    request: Request,
    principal: Principal = Depends(_principal),
) -> Principal:
    permission = (
        "documents.convert"
        if request.url.path.endswith("/convert")
        else "documents.edit"
    )
    _authorize(principal, permission, group="upload")
    return principal


class InputMessage(BaseModel):
    role: str
    text: str


class CreateTaskRequest(BaseModel):
    prompt: str = Field(min_length=1)
    file_ids: list[str] = Field(min_length=1)
    instructions: str = ""
    messages: list[InputMessage] = Field(default_factory=list)
    model: str | None = None
    client_request_id: str | None = Field(default=None, max_length=200)


class UpdateGatewayUserRequest(BaseModel):
    display_name: str | None = Field(default=None, max_length=200)
    enabled: bool | None = None
    is_admin: bool | None = None
    permissions: list[str] | None = None
    max_concurrent_tasks: int | None = Field(default=None, ge=1, le=100)
    storage_quota_bytes: int | None = Field(
        default=None, ge=MAX_FILE_BYTES, le=10 * 1024 * 1024 * 1024 * 1024
    )


def _task_row(task_id: str, owner_sub: str | None = None) -> sqlite3.Row:
    with _db() as db:
        if owner_sub is None:
            row = db.execute("SELECT * FROM tasks WHERE id = ?", (task_id,)).fetchone()
        else:
            row = db.execute(
                "SELECT * FROM tasks WHERE id = ? AND owner_sub = ?",
                (task_id, owner_sub),
            ).fetchone()
    if row is None:
        raise HTTPException(status_code=404, detail="任务不存在。")
    return row


def _task_json(row: sqlite3.Row) -> dict[str, Any]:
    with _db() as db:
        latest = db.execute(
            "SELECT COALESCE(MAX(id), 0) AS id FROM events WHERE task_id = ?",
            (row["id"],),
        ).fetchone()["id"]
    return {
        "id": row["id"],
        "status": row["status"],
        "output_text": row["output_text"],
        "progress": row["progress"],
        "detail": row["detail"],
        "error": row["error"],
        "last_event_id": latest,
        "created_at": row["created_at"],
        "updated_at": row["updated_at"],
    }


def _event(task_id: str, event_type: str, data: dict[str, Any]) -> int:
    with _db() as db:
        cursor = db.execute(
            "INSERT INTO events(task_id, type, data_json, created_at) VALUES (?, ?, ?, ?)",
            (task_id, event_type, json.dumps(data, ensure_ascii=False), _now()),
        )
        return int(cursor.lastrowid)


def _update_task(task_id: str, **values: Any) -> None:
    if not values:
        return
    values["updated_at"] = _now()
    columns = ", ".join(f"{key} = ?" for key in values)
    with _db() as db:
        db.execute(
            f"UPDATE tasks SET {columns} WHERE id = ?",  # noqa: S608; keys are internal constants
            (*values.values(), task_id),
        )


def _set_progress(task_id: str, progress: float, detail: str) -> None:
    bounded = max(0.0, min(1.0, progress))
    _update_task(task_id, status="running", progress=bounded, detail=detail)
    _event(task_id, "progress", {"progress": bounded, "detail": detail})


def _cancel_requested(task_id: str) -> bool:
    with _db() as db:
        row = db.execute(
            "SELECT cancel_requested FROM tasks WHERE id = ?", (task_id,)
        ).fetchone()
    return row is None or bool(row["cancel_requested"])


def _read_text_file(path: Path) -> str:
    raw = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _extract_text(path: Path, original_name: str) -> str:
    suffix = Path(original_name).suffix.lower()
    if suffix == ".pdf":
        return "\n\n".join(page.extract_text() or "" for page in PdfReader(path).pages)
    if suffix == ".docx":
        document = Document(path)
        blocks = [paragraph.text for paragraph in document.paragraphs]
        for table in document.tables:
            blocks.extend("\t".join(cell.text for cell in row.cells) for row in table.rows)
        return "\n".join(blocks)
    if suffix == ".xlsx":
        workbook = load_workbook(path, read_only=True, data_only=True)
        blocks: list[str] = []
        try:
            for sheet in workbook.worksheets:
                blocks.append(f"## Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    blocks.append("\t".join("" if value is None else str(value) for value in row))
        finally:
            workbook.close()
        return "\n".join(blocks)
    if suffix == ".pptx":
        deck = Presentation(path)
        blocks = []
        for index, slide in enumerate(deck.slides, start=1):
            blocks.append(f"## Slide {index}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    blocks.append(shape.text)
        return "\n".join(blocks)
    return _read_text_file(path)


def _chunks(text: str, size: int = CHUNK_CHARS) -> list[str]:
    normalized = re.sub(r"\r\n?", "\n", text).strip()
    if not normalized:
        return []
    chunks: list[str] = []
    cursor = 0
    while cursor < len(normalized):
        end = min(len(normalized), cursor + size)
        if end < len(normalized):
            boundary = normalized.rfind("\n", cursor + size // 2, end)
            if boundary > cursor:
                end = boundary
        chunks.append(normalized[cursor:end].strip())
        cursor = end
        while cursor < len(normalized) and normalized[cursor] == "\n":
            cursor += 1
    return [chunk for chunk in chunks if chunk]


def _chat_url() -> str:
    if not LLM_BASE_URL:
        raise RuntimeError("服务器未配置 LLM_BASE_URL。")
    if LLM_BASE_URL.endswith("/chat/completions"):
        return LLM_BASE_URL
    return f"{LLM_BASE_URL}/chat/completions"


def _headers() -> dict[str, str]:
    headers = {"Content-Type": "application/json"}
    if LLM_API_KEY:
        headers["Authorization"] = f"Bearer {LLM_API_KEY}"
    return headers


async def _completion(model: str, messages: list[dict[str, str]]) -> str:
    async with httpx.AsyncClient(timeout=httpx.Timeout(300.0, connect=30.0)) as client:
        response = await client.post(
            _chat_url(),
            headers=_headers(),
            json={"model": model, "messages": messages, "stream": False},
        )
        response.raise_for_status()
        payload = response.json()
    try:
        return str(payload["choices"][0]["message"]["content"] or "").strip()
    except (KeyError, IndexError, TypeError) as error:
        raise RuntimeError("上游模型没有返回可识别的文本。") from error


async def _stream_completion(
    task_id: str, model: str, messages: list[dict[str, str]]
) -> str:
    output = ""
    pending = ""
    last_flush = time.monotonic()
    async with httpx.AsyncClient(timeout=httpx.Timeout(600.0, connect=30.0)) as client:
        async with client.stream(
            "POST",
            _chat_url(),
            headers=_headers(),
            json={"model": model, "messages": messages, "stream": True},
        ) as response:
            response.raise_for_status()
            async for line in response.aiter_lines():
                if _cancel_requested(task_id):
                    raise asyncio.CancelledError
                if not line.startswith("data:"):
                    continue
                data = line[5:].strip()
                if not data or data == "[DONE]":
                    continue
                try:
                    payload = json.loads(data)
                    delta = payload["choices"][0]["delta"].get("content") or ""
                except (json.JSONDecodeError, KeyError, IndexError, TypeError):
                    continue
                if not delta:
                    continue
                output += delta
                pending += delta
                now = time.monotonic()
                if len(pending) >= 256 or now - last_flush >= 0.5:
                    _update_task(task_id, output_text=output, progress=0.98, detail="正在生成结果")
                    _event(task_id, "output_delta", {"delta": pending})
                    pending = ""
                    last_flush = now
    if pending:
        _update_task(task_id, output_text=output, progress=0.99, detail="正在保存结果")
        _event(task_id, "output_delta", {"delta": pending})
    return output.strip()


async def _reduce_material(
    task_id: str,
    model: str,
    prompt: str,
    material: list[str],
) -> str:
    """Hierarchically reduce map outputs so extremely long files stay bounded."""
    round_number = 0
    target_chars = CHUNK_CHARS * 4
    current = material
    while sum(len(item) for item in current) > target_chars:
        round_number += 1
        groups: list[list[str]] = []
        group: list[str] = []
        group_chars = 0
        for item in current:
            if group and group_chars + len(item) > CHUNK_CHARS * 3:
                groups.append(group)
                group = []
                group_chars = 0
            group.append(item)
            group_chars += len(item)
        if group:
            groups.append(group)

        reduced: list[str] = []
        for index, batch in enumerate(groups):
            if _cancel_requested(task_id):
                raise asyncio.CancelledError
            progress = 0.76 + 0.14 * (index + 1) / max(1, len(groups))
            _set_progress(
                task_id,
                progress,
                f"正在合并超长文档材料 {index + 1}/{len(groups)}",
            )
            reduced.append(
                await _completion(
                    model,
                    [
                        {
                            "role": "system",
                            "content": "合并下列文档阅读笔记。去重但保留与用户目标有关的事实、数字、证据出处、冲突和不确定性。",
                        },
                        {
                            "role": "user",
                            "content": f"用户目标：{prompt}\n\n" + "\n\n".join(batch),
                        },
                    ],
                )
            )
        if len(reduced) >= len(current):
            return "\n\n".join(reduced)[:target_chars]
        current = reduced
    return "\n\n".join(current)


def _history_messages(row: sqlite3.Row) -> list[dict[str, str]]:
    history: list[dict[str, str]] = []
    for item in json.loads(row["messages_json"]):
        role = item.get("role")
        text = str(item.get("text") or "").strip()
        if role in {"user", "assistant"} and text:
            history.append({"role": role, "content": text})
    return history[-12:]


async def _process_task(task_id: str) -> None:
    async with _semaphore:
        try:
            row = _task_row(task_id)
            if _cancel_requested(task_id):
                raise asyncio.CancelledError
            model = row["model"] or LLM_MODEL
            if not model:
                raise RuntimeError("服务器未配置 LLM_MODEL，任务也没有指定 model。")

            _set_progress(task_id, 0.03, "正在解析文件")
            file_ids: list[str] = json.loads(row["file_ids_json"])
            documents: list[tuple[str, str]] = []
            with _db() as db:
                for file_id in file_ids:
                    file_row = db.execute(
                        "SELECT * FROM files WHERE id = ?", (file_id,)
                    ).fetchone()
                    if file_row is None:
                        raise RuntimeError(f"文件 {file_id} 不存在或已被删除。")
                    text = await asyncio.to_thread(
                        _extract_text, Path(file_row["path"]), file_row["name"]
                    )
                    if text.strip():
                        documents.append((file_row["name"], text))
            if not documents:
                raise RuntimeError("上传文件中没有提取到可处理的文字。")

            all_chunks: list[tuple[str, str]] = []
            for name, text in documents:
                all_chunks.extend((name, chunk) for chunk in _chunks(text))
            if not all_chunks:
                raise RuntimeError("文档分段结果为空。")

            prompt = row["prompt"]
            summaries: list[str] = []
            for index, (name, chunk) in enumerate(all_chunks):
                if _cancel_requested(task_id):
                    raise asyncio.CancelledError
                progress = 0.08 + 0.62 * index / max(1, len(all_chunks))
                _set_progress(
                    task_id,
                    progress,
                    f"正在阅读文档片段 {index + 1}/{len(all_chunks)}",
                )
                summary = await _completion(
                    model,
                    [
                        {
                            "role": "system",
                            "content": "你负责长文档分析的分段阅读。只保留与用户目标有关的事实、证据、数字、结论和不确定性，不要编造。",
                        },
                        {
                            "role": "user",
                            "content": f"用户目标：{prompt}\n\n文件：{name}\n片段 {index + 1}/{len(all_chunks)}：\n{chunk}",
                        },
                    ],
                )
                summaries.append(f"### {name} / 片段 {index + 1}\n{summary}")

            _set_progress(task_id, 0.75, "已完成全文阅读，正在组织回答")
            reduced_material = await _reduce_material(
                task_id, model, prompt, summaries
            )
            instructions = row["instructions"].strip()
            final_messages: list[dict[str, str]] = [
                {
                    "role": "system",
                    "content": instructions
                    or "你正在完成长时间文档处理任务。严格依据材料，给出结构清晰、证据充分、可执行的最终结果。",
                },
                *_history_messages(row),
                {
                    "role": "user",
                    "content": f"用户当前目标：{prompt}\n\n以下是完整文档逐段阅读后的材料：\n\n"
                    + reduced_material,
                },
            ]
            output = await _stream_completion(task_id, model, final_messages)
            if not output:
                raise RuntimeError("上游模型完成了请求，但没有生成文本。")
            _update_task(
                task_id,
                status="completed",
                output_text=output,
                progress=1.0,
                detail="处理完成",
                error=None,
            )
            _event(task_id, "completed", {"output_text": output})
        except asyncio.CancelledError:
            if _shutting_down and not _cancel_requested(task_id):
                _update_task(
                    task_id,
                    status="queued",
                    detail="服务已停止，等待重启后恢复",
                    error=None,
                )
                _event(task_id, "interrupted", {"reason": "gateway_shutdown"})
            else:
                _update_task(
                    task_id,
                    status="cancelled",
                    detail="已取消",
                    error=None,
                    cancel_requested=1,
                )
                _event(task_id, "cancelled", {})
        except Exception as error:  # Persist failures so clients can retry later.
            message = str(error).strip() or error.__class__.__name__
            _update_task(task_id, status="failed", detail="处理失败", error=message)
            _event(task_id, "failed", {"error": message})
        finally:
            _running.pop(task_id, None)


def _schedule(task_id: str) -> None:
    if task_id in _running:
        return
    _running[task_id] = asyncio.create_task(_process_task(task_id))


async def _maintenance_loop() -> None:
    while True:
        await asyncio.sleep(CLEANUP_INTERVAL_SECONDS)
        await asyncio.to_thread(_cleanup_expired_tasks)


@asynccontextmanager
async def _lifespan(_: FastAPI):
    global _maintenance_task, _shutting_down
    _shutting_down = False
    _init_db()
    _cleanup_expired_tasks()
    with _db() as db:
        db.execute(
            "UPDATE tasks SET status = 'queued', detail = '服务已恢复，任务重新排队', updated_at = ? "
            "WHERE status IN ('queued', 'running') AND cancel_requested = 0",
            (_now(),),
        )
        task_ids = [
            row["id"]
            for row in db.execute(
                "SELECT id FROM tasks WHERE status = 'queued' AND cancel_requested = 0"
            ).fetchall()
        ]
    for task_id in task_ids:
        _schedule(task_id)
    if TASK_RETENTION_DAYS > 0:
        _maintenance_task = asyncio.create_task(_maintenance_loop())
    yield
    _shutting_down = True
    if _maintenance_task is not None:
        _maintenance_task.cancel()
        await asyncio.gather(_maintenance_task, return_exceptions=True)
        _maintenance_task = None
    for task in list(_running.values()):
        task.cancel()
    if _running:
        await asyncio.gather(*_running.values(), return_exceptions=True)


app = FastAPI(title="Expert Chat Gateway", version="0.3.1", lifespan=_lifespan)
app.dependency_overrides[document_require_auth] = _document_auth


@app.middleware("http")
async def security_headers(request: Request, call_next):
    response = await call_next(request)
    response.headers["X-Content-Type-Options"] = "nosniff"
    response.headers["X-Frame-Options"] = "DENY"
    response.headers["Referrer-Policy"] = "no-referrer"
    response.headers["Permissions-Policy"] = "camera=(), microphone=(), geolocation=()"
    if request.url.path == "/admin":
        response.headers["Content-Security-Policy"] = (
            "default-src 'none'; script-src 'unsafe-inline'; "
            "style-src 'unsafe-inline'; connect-src 'self'; "
            "base-uri 'none'; form-action 'self'; frame-ancestors 'none'"
        )
    if request.url.path.startswith("/v1/") or request.url.path == "/admin":
        response.headers["Cache-Control"] = "no-store"
    return response


@app.exception_handler(HTTPException)
async def http_exception_handler(_request: Request, exc: HTTPException) -> JSONResponse:
    # Document capabilities expose a structured error contract. Preserve it
    # when mounted here; retain FastAPI's normal {"detail": ...} shape for all
    # other Gateway endpoints.
    if isinstance(exc.detail, dict) and "error" in exc.detail:
        return JSONResponse(
            status_code=exc.status_code,
            content=exc.detail,
            headers=exc.headers,
        )
    return JSONResponse(
        status_code=exc.status_code,
        content={"detail": exc.detail},
        headers=exc.headers,
    )


@app.get("/health")
def health() -> dict[str, Any]:
    return {
        "ok": True,
        "version": app.version,
        "upstream_configured": bool(LLM_BASE_URL and LLM_MODEL),
        "model": LLM_MODEL or None,
        "active_tasks": len(_running),
        "authentication": AUTH.status(),
    }


@app.get("/v1/health")
def versioned_health() -> dict[str, Any]:
    return health()


@app.get("/v1/capabilities")
def capabilities(principal: Principal = Depends(_require("gateway.use"))) -> dict[str, Any]:
    entitlement = _provision_principal(principal)
    permissions = _permissions(entitlement)
    is_admin = bool(entitlement["is_admin"])
    manifest = MODULES.manifest()
    visible = {
        capability_id: metadata
        for capability_id, metadata in manifest.items()
        if is_admin
        or CAPABILITY_PERMISSIONS.get(capability_id, frozenset()).issubset(
            permissions
        )
    }
    return {
        "protocol_version": 1,
        "gateway_version": app.version,
        "capabilities": visible,
        "account": {
            "sub": principal.subject,
            "display_name": entitlement["display_name"],
            "is_admin": is_admin,
        },
    }


@app.get("/v1/me")
def current_account(principal: Principal = Depends(_principal)) -> dict[str, Any]:
    entitlement = _provision_principal(principal)
    with _db() as db:
        storage_used = db.execute(
            "SELECT COALESCE(SUM(size_bytes), 0) AS total FROM files WHERE owner_sub = ?",
            (principal.subject,),
        ).fetchone()["total"]
        active_tasks = db.execute(
            "SELECT COUNT(*) AS total FROM tasks WHERE owner_sub = ? "
            "AND status IN ('queued', 'running')",
            (principal.subject,),
        ).fetchone()["total"]
    return {
        "sub": principal.subject,
        "display_name": entitlement["display_name"],
        "auth_kind": principal.auth_kind,
        "enabled": bool(entitlement["enabled"]),
        "is_admin": bool(entitlement["is_admin"]),
        "permissions": sorted(_permissions(entitlement)),
        "quotas": {
            "storage_bytes": entitlement["storage_quota_bytes"],
            "storage_used_bytes": storage_used,
            "concurrent_tasks": entitlement["max_concurrent_tasks"],
            "active_tasks": active_tasks,
        },
    }


@app.post("/files", include_in_schema=False)
@app.post("/v1/files")
async def upload_file(
    file: UploadFile = File(...),
    principal: Principal = Depends(_require("files.write", group="upload")),
) -> dict[str, Any]:
    file_id = f"file_{uuid.uuid4().hex}"
    safe_suffix = Path(file.filename or "upload.bin").suffix[:16]
    path = UPLOAD_DIR / f"{file_id}{safe_suffix}"
    size = 0
    try:
        with path.open("wb") as target:
            while chunk := await file.read(1024 * 1024):
                size += len(chunk)
                if size > MAX_FILE_BYTES:
                    raise HTTPException(status_code=413, detail="文件超过 Gateway 单文件上限。")
                target.write(chunk)
    except Exception:
        path.unlink(missing_ok=True)
        raise
    entitlement = _provision_principal(principal)
    with _db() as db:
        used = db.execute(
            "SELECT COALESCE(SUM(size_bytes), 0) AS total FROM files WHERE owner_sub = ?",
            (principal.subject,),
        ).fetchone()["total"]
    if used + size > entitlement["storage_quota_bytes"]:
        path.unlink(missing_ok=True)
        raise HTTPException(status_code=413, detail="账户文件存储配额不足。")
    with _db() as db:
        db.execute(
            "INSERT INTO files(id, owner_sub, name, mime_type, size_bytes, path, created_at) "
            "VALUES (?, ?, ?, ?, ?, ?, ?)",
            (
                file_id,
                principal.subject,
                file.filename or "upload.bin",
                file.content_type or "application/octet-stream",
                size,
                str(path),
                _now(),
            ),
        )
    _audit(
        principal.subject,
        "file.upload",
        resource_type="file",
        resource_id=file_id,
        detail={"name": file.filename or "upload.bin", "size_bytes": size},
    )
    return {"id": file_id, "name": file.filename, "size_bytes": size}


@app.delete(
    "/files/{file_id}", include_in_schema=False
)
@app.delete("/v1/files/{file_id}")
def delete_file(
    file_id: str,
    principal: Principal = Depends(_require("files.write")),
) -> dict[str, Any]:
    with _db() as db:
        row = db.execute(
            "SELECT path FROM files WHERE id = ? AND owner_sub = ?",
            (file_id, principal.subject),
        ).fetchone()
        if row is None:
            raise HTTPException(status_code=404, detail="文件不存在。")
        active = db.execute(
            "SELECT 1 FROM tasks WHERE owner_sub = ? "
            "AND status IN ('queued', 'running') AND file_ids_json LIKE ? LIMIT 1",
            (principal.subject, f'%"{file_id}"%'),
        ).fetchone()
        if active is not None:
            raise HTTPException(status_code=409, detail="文件仍被运行中的任务使用。")
        db.execute(
            "DELETE FROM files WHERE id = ? AND owner_sub = ?",
            (file_id, principal.subject),
        )
    Path(row["path"]).unlink(missing_ok=True)
    _audit(
        principal.subject,
        "file.delete",
        resource_type="file",
        resource_id=file_id,
    )
    return {"deleted": True, "id": file_id}


@app.post("/tasks", include_in_schema=False)
@app.post("/v1/tasks")
async def create_task(
    request: CreateTaskRequest,
    principal: Principal = Depends(_require("tasks.create")),
) -> dict[str, Any]:
    entitlement = _provision_principal(principal)
    with _db() as db:
        client_request_id = (request.client_request_id or "").strip() or None
        if client_request_id is not None:
            existing_task = db.execute(
                "SELECT * FROM tasks WHERE owner_sub = ? AND client_request_id = ?",
                (principal.subject, client_request_id),
            ).fetchone()
            if existing_task is not None:
                return _task_json(existing_task)
        active_count = db.execute(
            "SELECT COUNT(*) AS total FROM tasks WHERE owner_sub = ? "
            "AND status IN ('queued', 'running')",
            (principal.subject,),
        ).fetchone()["total"]
        if active_count >= entitlement["max_concurrent_tasks"]:
            raise HTTPException(
                status_code=429,
                detail="账户同时运行的长任务已达到配额。",
                headers={"Retry-After": "15"},
            )
        existing = {
            row["id"]
            for row in db.execute(
                f"SELECT id FROM files WHERE owner_sub = ? AND id IN "
                f"({','.join('?' for _ in request.file_ids)})",
                (principal.subject, *request.file_ids),
            ).fetchall()
        }
        missing = [file_id for file_id in request.file_ids if file_id not in existing]
        if missing:
            raise HTTPException(status_code=400, detail=f"文件不存在：{', '.join(missing)}")
        task_id = f"task_{uuid.uuid4().hex}"
        now = _now()
        db.execute(
            """
            INSERT INTO tasks(
                id, owner_sub, client_request_id, status, prompt, instructions, messages_json, file_ids_json,
                model, output_text, progress, detail, created_at, updated_at
            ) VALUES (?, ?, ?, 'queued', ?, ?, ?, ?, ?, '', 0, '已排队', ?, ?)
            """,
            (
                task_id,
                principal.subject,
                client_request_id,
                request.prompt.strip(),
                request.instructions.strip(),
                json.dumps([item.model_dump() for item in request.messages], ensure_ascii=False),
                json.dumps(request.file_ids),
                (request.model or LLM_MODEL).strip(),
                now,
                now,
            ),
        )
    _event(task_id, "queued", {"detail": "已排队"})
    _schedule(task_id)
    _audit(
        principal.subject,
        "task.create",
        resource_type="task",
        resource_id=task_id,
        detail={"file_count": len(request.file_ids)},
    )
    return _task_json(_task_row(task_id, principal.subject))


@app.get(
    "/tasks/{task_id}", include_in_schema=False
)
@app.get("/v1/tasks/{task_id}")
def get_task(
    task_id: str,
    principal: Principal = Depends(_require("tasks.read")),
) -> dict[str, Any]:
    return _task_json(_task_row(task_id, principal.subject))


@app.get(
    "/tasks/{task_id}/events",
    include_in_schema=False,
)
@app.get("/v1/tasks/{task_id}/events")
def get_events(
    task_id: str,
    after: int = Query(default=0, ge=0),
    limit: int = Query(default=200, ge=1, le=1000),
    principal: Principal = Depends(_require("tasks.read")),
) -> dict[str, Any]:
    _task_row(task_id, principal.subject)
    with _db() as db:
        rows = db.execute(
            "SELECT * FROM events WHERE task_id = ? AND id > ? ORDER BY id LIMIT ?",
            (task_id, after, limit),
        ).fetchall()
    events = [
        {
            "id": row["id"],
            "type": row["type"],
            "data": json.loads(row["data_json"]),
            "created_at": row["created_at"],
        }
        for row in rows
    ]
    return {"events": events, "next_after": events[-1]["id"] if events else after}


@app.post(
    "/tasks/{task_id}/cancel",
    include_in_schema=False,
)
@app.post("/v1/tasks/{task_id}/cancel")
def cancel_task(
    task_id: str,
    principal: Principal = Depends(_require("tasks.cancel")),
) -> dict[str, Any]:
    row = _task_row(task_id, principal.subject)
    if row["status"] in {"completed", "failed", "cancelled"}:
        raise HTTPException(status_code=409, detail="任务已经结束。")
    _update_task(task_id, cancel_requested=1, detail="正在取消")
    _event(task_id, "cancel_requested", {})
    _audit(
        principal.subject,
        "task.cancel",
        resource_type="task",
        resource_id=task_id,
    )
    return _task_json(_task_row(task_id, principal.subject))


def _admin_user_json(row: sqlite3.Row) -> dict[str, Any]:
    return {
        "sub": row["owner_sub"],
        "display_name": row["display_name"],
        "enabled": bool(row["enabled"]),
        "is_admin": bool(row["is_admin"]),
        "permissions": sorted(_permissions(row)),
        "max_concurrent_tasks": row["max_concurrent_tasks"],
        "storage_quota_bytes": row["storage_quota_bytes"],
        "storage_used_bytes": row["storage_used_bytes"]
        if "storage_used_bytes" in row.keys()
        else 0,
        "active_tasks": row["active_tasks"] if "active_tasks" in row.keys() else 0,
        "created_at": row["created_at"],
        "updated_at": row["updated_at"],
    }


@app.get("/admin", response_class=FileResponse, include_in_schema=False)
def admin_console() -> FileResponse:
    return FileResponse(Path(__file__).with_name("admin.html"))


@app.get("/v1/admin/overview")
def admin_overview(
    principal: Principal = Depends(_require_admin),
) -> dict[str, Any]:
    del principal
    with _db() as db:
        users = db.execute("SELECT COUNT(*) AS total FROM user_entitlements").fetchone()[
            "total"
        ]
        storage = db.execute(
            "SELECT COALESCE(SUM(size_bytes), 0) AS total FROM files"
        ).fetchone()["total"]
        statuses = {
            row["status"]: row["total"]
            for row in db.execute(
                "SELECT status, COUNT(*) AS total FROM tasks GROUP BY status"
            ).fetchall()
        }
        recent_failures = [
            dict(row)
            for row in db.execute(
                "SELECT id, owner_sub, error, updated_at FROM tasks "
                "WHERE status = 'failed' ORDER BY updated_at DESC LIMIT 10"
            ).fetchall()
        ]
    return {
        "version": app.version,
        "users": users,
        "storage_bytes": storage,
        "active_workers": len(_running),
        "task_statuses": statuses,
        "recent_failures": recent_failures,
        "permission_catalog": PERMISSION_LABELS,
        "auth": AUTH.status(),
        "rate_limits": {
            "requests_per_minute": RATE_LIMIT_REQUESTS,
            "uploads_per_minute": RATE_LIMIT_UPLOADS,
        },
    }


@app.get("/v1/admin/users")
def admin_users(
    limit: int = Query(default=100, ge=1, le=500),
    offset: int = Query(default=0, ge=0),
    _: Principal = Depends(_require_admin),
) -> dict[str, Any]:
    with _db() as db:
        total = db.execute("SELECT COUNT(*) AS total FROM user_entitlements").fetchone()[
            "total"
        ]
        rows = db.execute(
            """
            SELECT u.*,
                COALESCE((SELECT SUM(f.size_bytes) FROM files f
                          WHERE f.owner_sub = u.owner_sub), 0) AS storage_used_bytes,
                COALESCE((SELECT COUNT(*) FROM tasks t
                          WHERE t.owner_sub = u.owner_sub
                          AND t.status IN ('queued', 'running')), 0) AS active_tasks
            FROM user_entitlements u
            ORDER BY u.updated_at DESC LIMIT ? OFFSET ?
            """,
            (limit, offset),
        ).fetchall()
    return {"items": [_admin_user_json(row) for row in rows], "total": total}


@app.put("/v1/admin/users/{owner_sub}")
def update_admin_user(
    owner_sub: str,
    request: UpdateGatewayUserRequest,
    principal: Principal = Depends(_require_admin),
) -> dict[str, Any]:
    if owner_sub == principal.subject and request.enabled is False:
        raise HTTPException(status_code=409, detail="不能停用当前管理员账户。")
    if owner_sub == principal.subject and request.is_admin is False:
        raise HTTPException(status_code=409, detail="不能移除当前账户的管理员权限。")
    updates: dict[str, Any] = {}
    if request.display_name is not None:
        name = request.display_name.strip()
        if not name:
            raise HTTPException(status_code=400, detail="显示名称不能为空。")
        updates["display_name"] = name
    if request.enabled is not None:
        updates["enabled"] = int(request.enabled)
    if request.is_admin is not None:
        updates["is_admin"] = int(request.is_admin)
    if request.permissions is not None:
        unknown = sorted(set(request.permissions) - set(PERMISSION_LABELS))
        if unknown:
            raise HTTPException(
                status_code=400, detail=f"未知权限：{', '.join(unknown)}"
            )
        updates["permissions_json"] = json.dumps(sorted(set(request.permissions)))
    if request.max_concurrent_tasks is not None:
        updates["max_concurrent_tasks"] = request.max_concurrent_tasks
    if request.storage_quota_bytes is not None:
        updates["storage_quota_bytes"] = request.storage_quota_bytes
    if not updates:
        raise HTTPException(status_code=400, detail="没有需要更新的字段。")
    updates["updated_at"] = _now()
    with _db() as db:
        existing = db.execute(
            "SELECT 1 FROM user_entitlements WHERE owner_sub = ?", (owner_sub,)
        ).fetchone()
        if existing is None:
            raise HTTPException(status_code=404, detail="Gateway 用户不存在。")
        assignments = ", ".join(f"{key} = ?" for key in updates)
        db.execute(
            f"UPDATE user_entitlements SET {assignments} WHERE owner_sub = ?",
            (*updates.values(), owner_sub),
        )
        row = db.execute(
            """
            SELECT u.*,
                COALESCE((SELECT SUM(f.size_bytes) FROM files f
                          WHERE f.owner_sub = u.owner_sub), 0) AS storage_used_bytes,
                COALESCE((SELECT COUNT(*) FROM tasks t
                          WHERE t.owner_sub = u.owner_sub
                          AND t.status IN ('queued', 'running')), 0) AS active_tasks
            FROM user_entitlements u WHERE u.owner_sub = ?
            """,
            (owner_sub,),
        ).fetchone()
    _audit(
        principal.subject,
        "user.update",
        target_sub=owner_sub,
        resource_type="user",
        resource_id=owner_sub,
        detail={"fields": sorted(updates)},
    )
    assert row is not None
    return _admin_user_json(row)


@app.get("/v1/admin/tasks")
def admin_tasks(
    status: str | None = Query(default=None),
    owner_sub: str | None = Query(default=None),
    limit: int = Query(default=100, ge=1, le=500),
    _: Principal = Depends(_require_admin),
) -> dict[str, Any]:
    clauses: list[str] = []
    values: list[Any] = []
    if status:
        clauses.append("status = ?")
        values.append(status)
    if owner_sub:
        clauses.append("owner_sub = ?")
        values.append(owner_sub)
    where = f"WHERE {' AND '.join(clauses)}" if clauses else ""
    with _db() as db:
        rows = db.execute(
            f"SELECT id, owner_sub, status, prompt, model, progress, detail, "
            f"error, created_at, updated_at FROM tasks {where} "
            f"ORDER BY updated_at DESC LIMIT ?",
            (*values, limit),
        ).fetchall()
    return {"items": [dict(row) for row in rows]}


@app.post("/v1/admin/tasks/{task_id}/cancel")
def admin_cancel_task(
    task_id: str,
    principal: Principal = Depends(_require_admin),
) -> dict[str, Any]:
    row = _task_row(task_id)
    if row["status"] in {"completed", "failed", "cancelled"}:
        raise HTTPException(status_code=409, detail="任务已经结束。")
    _update_task(task_id, cancel_requested=1, detail="管理员正在取消")
    _event(task_id, "cancel_requested", {"by": "admin"})
    _audit(
        principal.subject,
        "task.admin_cancel",
        target_sub=row["owner_sub"],
        resource_type="task",
        resource_id=task_id,
    )
    return _task_json(_task_row(task_id))


@app.get("/v1/admin/audit")
def admin_audit(
    limit: int = Query(default=100, ge=1, le=500),
    after_id: int = Query(default=0, ge=0),
    _: Principal = Depends(_require_admin),
) -> dict[str, Any]:
    with _db() as db:
        rows = db.execute(
            "SELECT * FROM audit_logs WHERE id > ? ORDER BY id DESC LIMIT ?",
            (after_id, limit),
        ).fetchall()
    items = []
    for row in rows:
        item = dict(row)
        item["detail"] = json.loads(item.pop("detail_json"))
        items.append(item)
    return {"items": items}


# Business modules are mounted through the same registry that publishes the
# capability manifest, so a route cannot silently diverge from discovery.
MODULES.install(app)
