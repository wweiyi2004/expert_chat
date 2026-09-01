"""Single-user document MCP Server with persistent resources and chunk upload."""

from __future__ import annotations

import base64
import binascii
import hmac
import os
import sqlite3
import uuid
from contextlib import contextmanager
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterator
from urllib.parse import urlsplit, urlunsplit

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from doc_edit.app.main import (
    DocumentOperationError,
    DocumentOutput,
    edit_document_path,
    convert_document_path,
)
from gateway.app.document_mcp import (
    DocumentMcpBackend,
    McpIdentity,
    StoredDocument,
    create_document_mcp_server,
    mcp_transport_security_from_env,
)


SERVER_VERSION = "1.0.0"
DATA_DIR = Path(os.getenv("MCP_DATA_DIR", "./mcp_data")).resolve()
DOCUMENT_DIR = DATA_DIR / "documents"
UPLOAD_DIR = DATA_DIR / "uploads"
DB_PATH = DATA_DIR / "documents.sqlite3"
MAX_FILE_BYTES = max(1, int(os.getenv("MCP_MAX_FILE_MB", "50"))) * 1024 * 1024
RESOURCE_MAX_CHARS = max(
    1_000,
    int(os.getenv("MCP_RESOURCE_MAX_CHARS", "200000")),
)
API_TOKEN = os.getenv("MCP_API_TOKEN", "").strip()
PUBLIC_URL = (
    os.getenv("MCP_PUBLIC_URL", "http://127.0.0.1:8790/mcp").strip()
    or "http://127.0.0.1:8790/mcp"
)

_SUPPORTED_UPLOADS = frozenset(
    {".pdf", ".xlsx", ".docx", ".pptx", ".txt", ".md", ".csv", ".tsv", ".json"}
)
_UPLOAD_CHUNK_BYTES = 1024 * 1024


def _issuer_url(public_url: str) -> str:
    parsed = urlsplit(public_url)
    return urlunsplit((parsed.scheme, parsed.netloc, "/", "", ""))


def _csv_env(name: str, default: str = "") -> list[str]:
    raw = os.getenv(name, default)
    return [item.strip() for item in raw.split(",") if item.strip()]


def _now() -> str:
    return datetime.now(timezone.utc).isoformat()


def _prepare_storage() -> None:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    DOCUMENT_DIR.mkdir(parents=True, exist_ok=True)
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    with _db() as db:
        db.executescript(
            """
            CREATE TABLE IF NOT EXISTS documents (
                id TEXT PRIMARY KEY,
                owner_sub TEXT NOT NULL,
                name TEXT NOT NULL,
                mime_type TEXT NOT NULL,
                size_bytes INTEGER NOT NULL,
                path TEXT NOT NULL,
                created_at TEXT NOT NULL
            );
            CREATE INDEX IF NOT EXISTS idx_documents_owner_created
                ON documents(owner_sub, created_at DESC);
            CREATE TABLE IF NOT EXISTS uploads (
                id TEXT PRIMARY KEY,
                owner_sub TEXT NOT NULL,
                name TEXT NOT NULL,
                mime_type TEXT NOT NULL,
                expected_size INTEGER NOT NULL,
                received_size INTEGER NOT NULL,
                path TEXT NOT NULL,
                created_at TEXT NOT NULL
            );
            """
        )


@contextmanager
def _db(*, immediate: bool = False) -> Iterator[sqlite3.Connection]:
    db = sqlite3.connect(DB_PATH, timeout=30)
    db.row_factory = sqlite3.Row
    try:
        db.execute("PRAGMA journal_mode=WAL")
        db.execute("BEGIN IMMEDIATE" if immediate else "BEGIN")
        yield db
        db.commit()
    except BaseException:
        db.rollback()
        raise
    finally:
        db.close()


def _safe_filename(filename: str) -> str:
    name = filename.strip()
    if (
        not name
        or name in {".", ".."}
        or "/" in name
        or "\\" in name
        or ".." in name
        or any(ord(char) < 32 for char in name)
    ):
        raise RuntimeError("文件名无效。")
    if len(name) > 180:
        raise RuntimeError("文件名过长。")
    suffix = Path(name).suffix.lower()
    if suffix not in _SUPPORTED_UPLOADS:
        raise RuntimeError(
            f"不支持 {suffix or '无扩展名'}；可用格式：{', '.join(sorted(_SUPPORTED_UPLOADS))}"
        )
    return name


def _stored(row: sqlite3.Row) -> StoredDocument:
    return StoredDocument(
        file_id=str(row["id"]),
        filename=str(row["name"]),
        mime_type=str(row["mime_type"]),
        size_bytes=int(row["size_bytes"]),
        created_at=str(row["created_at"]),
    )


def _checked_path(raw: str, root: Path) -> Path:
    path = Path(raw).resolve()
    try:
        path.relative_to(root.resolve())
    except ValueError as exc:
        raise RuntimeError("服务器存储路径无效。") from exc
    return path


def _document_row(owner_sub: str, file_id: str) -> sqlite3.Row:
    with _db() as db:
        row = db.execute(
            "SELECT * FROM documents WHERE id = ? AND owner_sub = ?",
            (file_id, owner_sub),
        ).fetchone()
    if row is None:
        raise RuntimeError("文档不存在。")
    path = _checked_path(str(row["path"]), DOCUMENT_DIR)
    if not path.is_file():
        raise RuntimeError("文档文件已丢失。")
    return row


def _resolve_token(token: str) -> McpIdentity | None:
    if not API_TOKEN:
        return None
    if not hmac.compare_digest(token.encode("utf-8"), API_TOKEN.encode("utf-8")):
        return None
    return McpIdentity(
        subject="single-user",
        display_name="Expert Chat User",
        client_id="expert-chat-app",
        scopes=(
            "gateway.use",
            "files.write",
            "documents.edit",
            "documents.convert",
        ),
        auth_kind="pre-shared-token",
    )


def _list_documents(owner_sub: str, limit: int) -> list[StoredDocument]:
    with _db() as db:
        rows = db.execute(
            "SELECT * FROM documents WHERE owner_sub = ? "
            "ORDER BY created_at DESC LIMIT ?",
            (owner_sub, limit),
        ).fetchall()
    return [_stored(row) for row in rows]


def _read_text_file(path: Path) -> str:
    raw = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _extract_text(path: Path, original_name: str) -> str:
    suffix = Path(original_name).suffix.lower()
    if suffix == ".pdf":
        return "\n\n".join(page.extract_text() or "" for page in PdfReader(path).pages)
    if suffix == ".docx":
        document = Document(path)
        blocks = [paragraph.text for paragraph in document.paragraphs]
        for table in document.tables:
            blocks.extend("\t".join(cell.text for cell in row.cells) for row in table.rows)
        return "\n".join(blocks)
    if suffix == ".xlsx":
        workbook = load_workbook(path, read_only=True, data_only=True)
        blocks: list[str] = []
        try:
            for sheet in workbook.worksheets:
                blocks.append(f"## Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    blocks.append("\t".join("" if value is None else str(value) for value in row))
        finally:
            workbook.close()
        return "\n".join(blocks)
    if suffix == ".pptx":
        deck = Presentation(path)
        blocks = []
        for index, slide in enumerate(deck.slides, start=1):
            blocks.append(f"## Slide {index}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    blocks.append(shape.text)
        return "\n".join(blocks)
    return _read_text_file(path)


def _metadata(owner_sub: str, file_id: str) -> dict[str, Any]:
    return _stored(_document_row(owner_sub, file_id)).to_dict()


def _inspect(owner_sub: str, file_id: str, max_chars: int) -> dict[str, Any]:
    row = _document_row(owner_sub, file_id)
    text = _extract_text(Path(row["path"]), str(row["name"]))
    preview = text[: min(max_chars, RESOURCE_MAX_CHARS)]
    result = _stored(row).to_dict()
    result.update(
        text=preview,
        text_chars=len(text),
        truncated=len(preview) < len(text),
    )
    return result


def _read_text(owner_sub: str, file_id: str) -> str:
    row = _document_row(owner_sub, file_id)
    text = _extract_text(Path(row["path"]), str(row["name"]))
    if len(text) <= RESOURCE_MAX_CHARS:
        return text
    return text[:RESOURCE_MAX_CHARS] + f"\n\n…（已截断；原文 {len(text)} 字符）"


def _read_binary(owner_sub: str, file_id: str) -> bytes:
    row = _document_row(owner_sub, file_id)
    return _checked_path(str(row["path"]), DOCUMENT_DIR).read_bytes()


def _begin_upload(
    owner_sub: str,
    filename: str,
    mime_type: str,
    size_bytes: int,
) -> dict[str, Any]:
    name = _safe_filename(filename)
    if size_bytes <= 0 or size_bytes > MAX_FILE_BYTES:
        raise RuntimeError(f"文件大小必须在 1 到 {MAX_FILE_BYTES} 字节之间。")
    upload_id = f"upload_{uuid.uuid4().hex}"
    path = UPLOAD_DIR / f"{upload_id}.part"
    path.touch(exist_ok=False)
    try:
        with _db(immediate=True) as db:
            db.execute(
                "INSERT INTO uploads(id, owner_sub, name, mime_type, expected_size, "
                "received_size, path, created_at) VALUES (?, ?, ?, ?, ?, 0, ?, ?)",
                (
                    upload_id,
                    owner_sub,
                    name,
                    mime_type.strip() or "application/octet-stream",
                    size_bytes,
                    str(path),
                    _now(),
                ),
            )
    except BaseException:
        path.unlink(missing_ok=True)
        raise
    return {
        "upload_id": upload_id,
        "offset": 0,
        "chunk_size_bytes": _UPLOAD_CHUNK_BYTES,
    }


def _append_upload(
    owner_sub: str,
    upload_id: str,
    chunk_base64: str,
    offset: int,
) -> dict[str, Any]:
    try:
        chunk = base64.b64decode(chunk_base64, validate=True)
    except (binascii.Error, ValueError) as exc:
        raise RuntimeError("上传分块不是合法 Base64。") from exc
    if not chunk or len(chunk) > _UPLOAD_CHUNK_BYTES:
        raise RuntimeError(f"上传分块必须在 1 到 {_UPLOAD_CHUNK_BYTES} 字节之间。")
    with _db(immediate=True) as db:
        row = db.execute(
            "SELECT * FROM uploads WHERE id = ? AND owner_sub = ?",
            (upload_id, owner_sub),
        ).fetchone()
        if row is None:
            raise RuntimeError("上传会话不存在。")
        received = int(row["received_size"])
        expected = int(row["expected_size"])
        if offset != received:
            raise RuntimeError(f"上传 offset 不匹配；服务器需要 {received}。")
        if received + len(chunk) > expected or received + len(chunk) > MAX_FILE_BYTES:
            raise RuntimeError("上传内容超过声明大小或服务器上限。")
        path = _checked_path(str(row["path"]), UPLOAD_DIR)
        with path.open("ab") as output:
            output.write(chunk)
        next_offset = received + len(chunk)
        db.execute(
            "UPDATE uploads SET received_size = ? WHERE id = ?",
            (next_offset, upload_id),
        )
    return {
        "upload_id": upload_id,
        "offset": next_offset,
        "complete": next_offset == expected,
    }


def _finish_upload(owner_sub: str, upload_id: str) -> StoredDocument:
    with _db(immediate=True) as db:
        row = db.execute(
            "SELECT * FROM uploads WHERE id = ? AND owner_sub = ?",
            (upload_id, owner_sub),
        ).fetchone()
        if row is None:
            raise RuntimeError("上传会话不存在。")
        expected = int(row["expected_size"])
        received = int(row["received_size"])
        if received != expected:
            raise RuntimeError(f"上传尚未完成：{received}/{expected} 字节。")
        source = _checked_path(str(row["path"]), UPLOAD_DIR)
        if not source.is_file() or source.stat().st_size != expected:
            raise RuntimeError("上传临时文件大小不一致。")
        file_id = f"file_{uuid.uuid4().hex}"
        suffix = Path(str(row["name"])).suffix.lower()
        target = DOCUMENT_DIR / f"{file_id}{suffix}"
        source.replace(target)
        created_at = _now()
        try:
            db.execute(
                "INSERT INTO documents(id, owner_sub, name, mime_type, size_bytes, path, created_at) "
                "VALUES (?, ?, ?, ?, ?, ?, ?)",
                (
                    file_id,
                    owner_sub,
                    str(row["name"]),
                    str(row["mime_type"]),
                    expected,
                    str(target),
                    created_at,
                ),
            )
            db.execute("DELETE FROM uploads WHERE id = ?", (upload_id,))
        except BaseException:
            target.unlink(missing_ok=True)
            raise
    return StoredDocument(
        file_id=file_id,
        filename=str(row["name"]),
        mime_type=str(row["mime_type"]),
        size_bytes=expected,
        created_at=created_at,
    )


def _abort_upload(owner_sub: str, upload_id: str) -> dict[str, Any]:
    with _db(immediate=True) as db:
        row = db.execute(
            "SELECT path FROM uploads WHERE id = ? AND owner_sub = ?",
            (upload_id, owner_sub),
        ).fetchone()
        if row is not None:
            db.execute("DELETE FROM uploads WHERE id = ?", (upload_id,))
    if row is not None:
        _checked_path(str(row["path"]), UPLOAD_DIR).unlink(missing_ok=True)
    return {"upload_id": upload_id, "deleted": row is not None}


def _store_output(
    owner_sub: str,
    file_id: str,
    path: Path,
    output: DocumentOutput,
) -> StoredDocument:
    size = path.stat().st_size
    if size <= 0 or size > MAX_FILE_BYTES:
        path.unlink(missing_ok=True)
        raise RuntimeError("生成文件为空或超过服务器限制。")
    created_at = _now()
    try:
        with _db(immediate=True) as db:
            db.execute(
                "INSERT INTO documents(id, owner_sub, name, mime_type, size_bytes, path, created_at) "
                "VALUES (?, ?, ?, ?, ?, ?, ?)",
                (
                    file_id,
                    owner_sub,
                    output.filename,
                    output.media_type,
                    size,
                    str(path),
                    created_at,
                ),
            )
    except BaseException:
        path.unlink(missing_ok=True)
        raise
    return StoredDocument(
        file_id=file_id,
        filename=output.filename,
        mime_type=output.media_type,
        size_bytes=size,
        created_at=created_at,
    )


def _edit(
    owner_sub: str,
    source_id: str,
    patch: dict[str, Any],
) -> StoredDocument:
    row = _document_row(owner_sub, source_id)
    file_id = f"file_{uuid.uuid4().hex}"
    suffix = Path(str(row["name"])).suffix.lower()
    path = DOCUMENT_DIR / f"{file_id}{suffix}"
    try:
        output = edit_document_path(
            Path(str(row["path"])),
            path,
            str(row["name"]),
            patch,
        )
        return _store_output(owner_sub, file_id, path, output)
    except DocumentOperationError as exc:
        path.unlink(missing_ok=True)
        raise RuntimeError(f"{exc.code}: {exc.message}") from exc
    except BaseException:
        path.unlink(missing_ok=True)
        raise


def _convert(
    owner_sub: str,
    source_id: str,
    target_format: str,
    output_filename: str | None,
) -> StoredDocument:
    row = _document_row(owner_sub, source_id)
    target = target_format.strip().lower().lstrip(".")
    file_id = f"file_{uuid.uuid4().hex}"
    path = DOCUMENT_DIR / f"{file_id}.{target}"
    try:
        output = convert_document_path(
            Path(str(row["path"])),
            path,
            str(row["name"]),
            target,
            output_filename,
        )
        return _store_output(owner_sub, file_id, path, output)
    except DocumentOperationError as exc:
        path.unlink(missing_ok=True)
        raise RuntimeError(f"{exc.code}: {exc.message}") from exc
    except BaseException:
        path.unlink(missing_ok=True)
        raise


def _delete(owner_sub: str, file_id: str) -> dict[str, Any]:
    with _db(immediate=True) as db:
        row = db.execute(
            "SELECT path FROM documents WHERE id = ? AND owner_sub = ?",
            (file_id, owner_sub),
        ).fetchone()
        if row is not None:
            db.execute("DELETE FROM documents WHERE id = ?", (file_id,))
    if row is not None:
        _checked_path(str(row["path"]), DOCUMENT_DIR).unlink(missing_ok=True)
    return {"file_id": file_id, "deleted": row is not None}


_prepare_storage()
DOCUMENT_MCP = create_document_mcp_server(
    DocumentMcpBackend(
        resolve_token=_resolve_token,
        list_documents=_list_documents,
        inspect_document=_inspect,
        read_document_metadata=_metadata,
        read_document_text=_read_text,
        read_document_binary=_read_binary,
        edit_document=_edit,
        convert_document=_convert,
        begin_upload=_begin_upload,
        append_upload=_append_upload,
        finish_upload=_finish_upload,
        abort_upload=_abort_upload,
        delete_document=_delete,
    ),
    issuer_url=os.getenv("MCP_ISSUER", "").strip() or _issuer_url(PUBLIC_URL),
    resource_server_url=PUBLIC_URL,
    version=SERVER_VERSION,
)

app = DOCUMENT_MCP.streamable_http_app(
    streamable_http_path="/mcp",
    json_response=True,
    stateless_http=True,
    max_request_body_size=2 * 1024 * 1024,
    transport_security=mcp_transport_security_from_env(
        allowed_hosts=_csv_env(
            "MCP_ALLOWED_HOSTS",
            "127.0.0.1:*,localhost:*,[::1]:*",
        ),
        allowed_origins=_csv_env("MCP_ALLOWED_ORIGINS"),
    ),
)
